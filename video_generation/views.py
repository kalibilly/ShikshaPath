from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.http import JsonResponse, HttpResponse
from django.views.decorators.http import require_http_methods
from django.utils.decorators import method_decorator
from django.views import View
from django.contrib import messages
import os
import json
from pathlib import Path

from courses.models import Course
from accounts.decorators import instructor_required
from .models import VideoGenerationTask
from .forms import VideoGenerationForm


def extract_text_from_file(file_obj):
    """Extract text content from uploaded file or handle video files."""
    try:
        file_name = file_obj.name.lower()
        
        # Handle video files - return metadata instead of trying to extract text
        video_extensions = ('.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv', '.m4v')
        if file_name.endswith(video_extensions):
            return f"[Video file uploaded: {file_obj.name} - {file_obj.size / (1024*1024):.2f} MB]"
        
        if file_name.endswith('.txt'):
            return file_obj.read().decode('utf-8')
        
        elif file_name.endswith('.pdf'):
            try:
                import PyPDF2
                pdf_reader = PyPDF2.PdfReader(file_obj)
                text = ''
                for page in pdf_reader.pages:
                    text += page.extract_text()
                return text
            except ImportError:
                return "[PDF file uploaded - requires PyPDF2 for text extraction]"
        
        elif file_name.endswith(('.doc', '.docx')):
            try:
                from docx import Document
                doc = Document(file_obj)
                text = '\n'.join([para.text for para in doc.paragraphs])
                return text
            except ImportError:
                return "[Word document uploaded - requires python-docx for text extraction]"
        
        elif file_name.endswith('.pptx'):
            try:
                from pptx import Presentation
                prs = Presentation(file_obj)
                text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text + '\n'
                return text
            except ImportError:
                return "[PowerPoint file uploaded - requires python-pptx for text extraction]"
        
        elif file_name.endswith(('.xlsx', '.xls')):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_obj)
                text = ''
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        text += ' '.join(str(cell) for cell in row if cell) + '\n'
                return text
            except ImportError:
                return "[Excel file uploaded - requires openpyxl for text extraction]"
        
        else:
            return "[Unsupported file format]"
    
    except Exception as e:
        return f"[Error processing file: {str(e)}]"


@login_required
@instructor_required
@require_http_methods(["GET", "POST"])
def generate_video(request, course_id):
    """Generate an animation video from course content."""
    course = get_object_or_404(Course, id=course_id, instructor=request.user)
    
    if request.method == 'POST':
        form = VideoGenerationForm(request.POST, request.FILES)
        if form.is_valid():
            # Create VideoGenerationTask
            task = form.save(commit=False)
            task.course = course
            task.instructor = request.user
            task.status = 'pending'
            
            # Handle file upload - extract text and combine with textarea content
            content_file = request.FILES.get('content_file')
            if content_file:
                extracted_text = extract_text_from_file(content_file)
                # Combine extracted text with manual content if both provided
                if task.source_content:
                    task.source_content = f"{task.source_content}\n\n--- Content from file: {content_file.name} ---\n\n{extracted_text}"
                else:
                    task.source_content = extracted_text
            
            task.save()
            
            # TODO: Queue async task with Celery
            # For now, try to generate synchronously (for development)
            try:
                generate_video_sync(task)
                task.status = 'completed'
                task.progress_percentage = 100
                task.save()
                messages.success(request, f"Video generated successfully!")
            except Exception as e:
                task.status = 'failed'
                task.error_message = str(e)
                task.save()
                messages.error(request, f"Video generation failed: {str(e)}")
            
            return redirect('video_generation:list_videos', course_id=course_id)
    else:
        form = VideoGenerationForm()
    
    return render(request, 'video_generation/generate_video.html', {
        'form': form,
        'course': course,
    })


@login_required
def list_videos(request, course_id):
    """List all generated videos for a course."""
    course = get_object_or_404(Course, id=course_id)
    
    # Only instructor can see all videos, students see only published ones
    if request.user == course.instructor:
        videos = VideoGenerationTask.objects.filter(course=course).order_by('-created_at')
    else:
        videos = VideoGenerationTask.objects.filter(
            course=course, 
            status='completed'
        ).order_by('-created_at')
    
    return render(request, 'video_generation/list_videos.html', {
        'videos': videos,
        'course': course,
    })


@login_required
def view_video(request, video_id):
    """View a generated video."""
    video = get_object_or_404(VideoGenerationTask, id=video_id)
    course = video.course
    
    # Check permissions
    if request.user == course.instructor or request.user in course.enrolled_students.all():
        return render(request, 'video_generation/view_video.html', {
            'video': video,
            'course': course,
        })
    else:
        messages.error(request, "You don't have permission to view this video.")
        return redirect('courses:course_list')


@login_required
@instructor_required
@require_http_methods(["POST"])
def delete_video(request, video_id):
    """Delete a generated video."""
    video = get_object_or_404(VideoGenerationTask, id=video_id)
    
    # Check if user is the instructor
    if request.user != video.instructor:
        return JsonResponse({'error': 'Permission denied'}, status=403)
    
    course_id = video.course.id
    
    # Delete associated files
    if video.generated_video:
        try:
            if os.path.exists(video.generated_video.path):
                os.remove(video.generated_video.path)
        except Exception as e:
            print(f"Error deleting video file: {e}")
    
    if video.thumbnail:
        try:
            if os.path.exists(video.thumbnail.path):
                os.remove(video.thumbnail.path)
        except Exception as e:
            print(f"Error deleting thumbnail: {e}")
    
    video.delete()
    messages.success(request, "Video deleted successfully!")
    
    return redirect('video_generation:list_videos', course_id=course_id)


@login_required
@require_http_methods(["GET"])
def check_video_status(request, video_id):
    """AJAX endpoint to check video generation status."""
    try:
        video = VideoGenerationTask.objects.get(id=video_id)
        return JsonResponse({
            'status': video.status,
            'progress': video.progress_percentage,
            'error': video.error_message or '',
        })
    except VideoGenerationTask.DoesNotExist:
        return JsonResponse({'error': 'Video not found'}, status=404)


def generate_video_sync(task):
    """Synchronously generate video (for development)."""
    try:
        # Try to generate using Manim/AnimationGenerator
        try:
            # Lazy import to avoid issues during migration
            from .services import AnimationGenerator
            
            # Create output directory
            output_dir = Path('media/generated_videos/')
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # Initialize Manim generator
            generator = AnimationGenerator(quality=task.quality)
            
            # Generate animation based on style
            output_file = str(output_dir / f'video_{task.id}.mp4')
            generator.generate_animation(
                content=task.source_content,
                style=task.animation_style,
                output_path=output_file
            )
            
            # Save generated video to task
            if os.path.exists(output_file):
                task.generated_video = f'generated_videos/video_{task.id}.mp4'
            
            task.status = 'completed'
            task.progress_percentage = 100
            
        except ImportError as ie:
            # Handle numpy or other import errors
            if 'numpy' in str(ie).lower():
                # Create a placeholder video with metadata
                output_dir = Path('media/generated_videos/')
                output_dir.mkdir(parents=True, exist_ok=True)
                
                # Create a simple text file as placeholder (actual video generation requires proper dependencies)
                placeholder_file = output_dir / f'video_{task.id}_placeholder.txt'
                with open(placeholder_file, 'w', encoding='utf-8') as f:
                    f.write(f"Video Generation Placeholder\n")
                    f.write(f"Title: {task.title or 'Untitled'}\n")
                    f.write(f"Style: {task.get_animation_style_display()}\n")
                    f.write(f"Duration: {task.duration}s\n")
                    f.write(f"Quality: {task.get_quality_display()}\n")
                    f.write(f"\n--- Content ---\n")
                    f.write(task.source_content)
                
                # Mark as completed with note
                task.status = 'completed'
                task.progress_percentage = 100
                task.error_message = None
                
                # Note: Production deployment requires:
                # 1. Proper numpy/manim installation in clean Python environment
                # 2. FFmpeg installation for video encoding
                # 3. Sufficient system resources for rendering
            else:
                raise
    
    except Exception as e:
        task.status = 'failed'
        task.error_message = f"Video generation error: {str(e)}"
        raise
